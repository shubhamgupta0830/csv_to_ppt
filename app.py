import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def add_question_slide(prs, question, options, correct=None, show_answer=False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add question
    title_shape = shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    frame = title_shape.text_frame
    p = frame.add_paragraph()
    p.text = f"Q: {question}"
    p.font.size = Pt(24)

    # Add options
    top = 1.5
    for opt in options:
        para = shapes.add_textbox(Inches(1), Inches(top), Inches(8), Inches(0.5)).text_frame.add_paragraph()
        para.text = f"- {opt}"
        para.font.size = Pt(18)
        if show_answer and str(opt).strip() == str(correct).strip():
            para.font.bold = True
        top += 0.6

def generate_ppt(df):
    prs = Presentation()

    # Detect option columns
    option_columns = [col for col in df.columns if col.startswith("option")]

    for _, row in df.iterrows():
        question = row['question']
        options = [row[col] for col in option_columns if pd.notna(row[col])]
        correct = row['correct']

        # Slide 1: Only question and options
        add_question_slide(prs, question, options, correct, show_answer=False)

        # Slide 2: Reveal correct answer
        add_question_slide(prs, question, options, correct, show_answer=True)

    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

# Streamlit app
st.title("MCQ to PowerPoint Generator (With Answer Reveal)")

uploaded_file = st.file_uploader("Upload CSV file with questions", type=["csv"])

if uploaded_file:
    df = pd.read_csv(uploaded_file)

    if 'question' not in df.columns or 'correct' not in df.columns:
        st.error("CSV must have 'question' and 'correct' columns.")
    else:
        ppt_bytes = generate_ppt(df)
        st.success("PowerPoint with answer reveal generated!")

        st.download_button(
            label="Download PowerPoint",
            data=ppt_bytes,
            file_name="mcq_questions_with_reveal.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )